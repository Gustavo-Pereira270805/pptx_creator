#!/usr/bin/env python3
"""PPTX Motion — inject entrance animations and transitions into a .pptx.

python-pptx has no public animation API, so this helper edits the slide XML
directly (p:timing for entrance effects, p:transition for slide transitions,
including Morph). See ../../references/pptx-motion.md for the choreography rules.

Usage:
  python3 pptx_motion.py deck.pptx --entrance <shape> <effect> [trigger] [delay_ms] --out out.pptx
  python3 pptx_motion.py deck.pptx --transition <fade|push|wipe|morph|none> [--slide N|all] --out out.pptx
  python3 pptx_motion.py deck.pptx --spec motion-spec.json --out out.pptx

Effects: appear | fade | zoom        Triggers: on_click | after_previous | with_previous
"""

import argparse
import json
import sys

from lxml import etree
from pptx import Presentation

P = "http://schemas.openxmlformats.org/presentationml/2006/main"
MC = "http://schemas.openxmlformats.org/markup-compatibility/2006"
P14 = "http://schemas.microsoft.com/office/powerpoint/2010/main"

EFFECTS = {
    "appear": {"preset_id": "1", "subtype": "0", "kind": "appear"},
    "fade": {"preset_id": "10", "subtype": "0", "kind": "fade"},
    "zoom": {"preset_id": "23", "subtype": "16", "kind": "zoom"},
}
TRIGGERS = {"on_click": "clickEffect", "after_previous": "afterEffect", "with_previous": "withEffect"}
TRANSITIONS = {"fade", "push", "wipe", "morph", "none"}


def _build_effect_block(shape_id, effect, trigger, delay_ms, id_counter):
    """Return (par_element, next_id) for one entrance effect."""
    spec = EFFECTS[effect]
    preset = spec["preset_id"]
    subtype = spec["subtype"]
    node_type = TRIGGERS[trigger]

    def nid():
        nonlocal id_counter
        id_counter += 1
        return str(id_counter)

    delay = str(delay_ms)
    outer_cond = "indefinite" if trigger == "on_click" else "0"

    xml = f"""<p:par xmlns:p="{P}">
      <p:cTn id="{nid()}" fill="hold">
        <p:stCondLst><p:cond delay="{outer_cond}"/></p:stCondLst>
        <p:childTnLst>
          <p:par>
            <p:cTn id="{nid()}" fill="hold">
              <p:stCondLst><p:cond delay="{delay}"/></p:stCondLst>
              <p:childTnLst>
                <p:par>
                  <p:cTn id="{nid()}" presetID="{preset}" presetClass="entr" presetSubtype="{subtype}" fill="hold" nodeType="{node_type}">
                    <p:stCondLst><p:cond delay="0"/></p:stCondLst>
                    <p:childTnLst>
                      <p:set>
                        <p:cBhvr>
                          <p:cTn id="{nid()}" dur="1" fill="hold">
                            <p:stCondLst><p:cond delay="0"/></p:stCondLst>
                          </p:cTn>
                          <p:tgtEl><p:spTgt spid="{shape_id}"/></p:tgtEl>
                          <p:attrNameLst><p:attrName>style.visibility</p:attrName></p:attrNameLst>
                        </p:cBhvr>
                        <p:to><p:strVal val="visible"/></p:to>
                      </p:set>"""
    if spec["kind"] == "fade":
        xml += f"""
                      <p:animEffect transition="in" filter="fade">
                        <p:cBhvr>
                          <p:cTn id="{nid()}" dur="500" fill="hold">
                            <p:stCondLst><p:cond delay="0"/></p:stCondLst>
                          </p:cTn>
                          <p:tgtEl><p:spTgt spid="{shape_id}"/></p:tgtEl>
                        </p:cBhvr>
                      </p:animEffect>"""
    elif spec["kind"] == "zoom":
        xml += f"""
                      <p:animScale>
                        <p:cBhvr>
                          <p:cTn id="{nid()}" dur="600" fill="hold">
                            <p:stCondLst><p:cond delay="0"/></p:stCondLst>
                          </p:cTn>
                          <p:tgtEl><p:spTgt spid="{shape_id}"/></p:tgtEl>
                        </p:cBhvr>
                        <p:from x="0" y="0"/>
                        <p:to x="100000" y="100000"/>
                      </p:animScale>"""
    xml += """
                    </p:childTnLst>
                  </p:cTn>
                </p:par>
              </p:childTnLst>
            </p:cTn>
          </p:par>
        </p:childTnLst>
      </p:cTn>
    </p:par>"""
    return etree.fromstring(xml), id_counter


def _get_timing(slide):
    sld = slide._element
    timing = sld.find(etree.QName(P, "timing"))
    if timing is None:
        # build the standard tmRoot skeleton
        timing = etree.fromstring(f"""<p:timing xmlns:p="{P}">
  <p:tnLst><p:par><p:cTn id="1" dur="indefinite" restart="never" nodeType="tmRoot">
    <p:childTnLst><p:seq concurrent="1" nextAc="seek">
      <p:cTn id="2" dur="indefinite" nodeType="mainSeq"><p:childTnLst/></p:cTn>
      <p:prevCondLst><p:cond evt="onPrev" delay="0"><p:tgtEl><p:sldTgt/></p:tgtEl></p:cond></p:prevCondLst>
      <p:nextCondLst><p:cond evt="onNext" delay="0"><p:tgtEl><p:sldTgt/></p:tgtEl></p:cond></p:nextCondLst>
    </p:seq></p:childTnLst>
  </p:cTn></p:par></p:tnLst>
</p:timing>""")
        sld.append(timing)
        # keep schema order: cSld, clrMapOvr, transition, timing
        transition = sld.find(etree.QName(P, "transition"))
        if transition is not None:
            sld.remove(timing)
            transition.addnext(timing)
    return timing


def _main_seq(timing):
    return timing.find(f".//{{{P}}}seq/{{{P}}}cTn/{{{P}}}childTnLst")


def _shape_by_name(slide, name):
    for shape in slide.shapes:
        if shape.name == name:
            return shape
    return None


def _shape_by_index(slide, index):
    shapes = list(slide.shapes)
    return shapes[index] if 0 <= index < len(shapes) else None


def add_entrance(slide, target, effect="fade", trigger="after_previous", delay_ms=0):
    if effect not in EFFECTS:
        raise ValueError(f"Efeito desconhecido: {effect} (use: {', '.join(EFFECTS)})")
    if trigger not in TRIGGERS:
        raise ValueError(f"Trigger desconhecido: {trigger} (use: {', '.join(TRIGGERS)})")

    if target.startswith("index:"):
        shape = _shape_by_index(slide, int(target.split(":", 1)[1]))
    else:
        shape = _shape_by_name(slide, target)
    if shape is None:
        return f"SKIP: shape '{target}' não encontrada no slide"
    shape_id = str(shape.shape_id)

    timing = _get_timing(slide)
    seq_child = _main_seq(timing)

    # reuse existing ids if any, else start above the skeleton ids
    used = [int(e.get("id")) for e in timing.iter(etree.QName(P, "cTn"))]
    counter = max(used) if used else 1

    block, counter = _build_effect_block(shape_id, effect, trigger, delay_ms, counter)
    seq_child.append(block)
    return f"OK: '{target}' (spid {shape_id}) <- {effect} ({trigger}{' +%dms' % delay_ms if delay_ms else ''})"


def add_transition(slide, kind="fade", duration_ms=None, direction=None):
    sld = slide._element
    # remove any existing transition (including mc:AlternateContent morph)
    for tag in ("transition",):
        for el in sld.iter(etree.QName(P, tag)):
            parent = el.getparent()
            if parent is not None:
                parent.remove(el)
    # AlternateContent morph lives as a direct child with its own namespace
    for el in list(sld):
        if isinstance(el.tag, str) and etree.QName(el.tag).namespace == MC:
            sld.remove(el)

    if kind == "none":
        return "OK: transição removida"

    if kind == "morph":
        xml = f"""<mc:AlternateContent xmlns:mc="{MC}" xmlns:p="{P}" xmlns:p14="{P14}">
      <mc:Choice Requires="p14">
        <p:transition spd="med" p14:dur="{duration_ms or 600}">
          <p14:morph>
            <p14:spTgt spid="0"/>
          </p14:morph>
        </p:transition>
      </mc:Choice>
      <mc:Fallback>
        <p:transition spd="med"><p:fade/></p:transition>
      </mc:Fallback>
    </mc:AlternateContent>"""
        el = etree.fromstring(xml)
    else:
        inner = {"fade": "<p:fade/>", "push": f'<p:push dir="{direction or "l"}"/>',
                 "wipe": f'<p:wipe dir="{direction or "l"}"/>'}[kind]
        dur = f' p14:dur="{duration_ms}"' if duration_ms else ""
        xml = f'<p:transition xmlns:p="{P}" spd="med"{dur}>{inner}</p:transition>'
        el = etree.fromstring(xml)

    # insert transition before timing (schema order: cSld, clrMapOvr, transition, timing)
    timing = sld.find(etree.QName(P, "timing"))
    if timing is not None:
        timing.addprevious(el)
    else:
        sld.append(el)
    return f"OK: transição '{kind}' no slide"


def apply_motion_spec(prs, spec):
    report = []
    transitions = spec.get("transitions", {})
    default_trans = transitions.get("default")

    for idx, slide in enumerate(prs.slides, start=1):
        t_kind = transitions.get(str(idx), default_trans)
        if t_kind:
            report.append(f"slide {idx}: {add_transition(slide, t_kind)}")

    for pair in spec.get("morph_pairs", []):
        for idx in pair:
            if 1 <= int(idx) <= len(prs.slides):
                report.append(f"slide {idx}: {add_transition(prs.slides[int(idx) - 1], 'morph')}")

    for idx_str, anims in (spec.get("slides") or {}).items():
        slide = prs.slides[int(idx_str) - 1]
        if isinstance(anims, dict) and "stagger" in anims:
            st = anims["stagger"]
            shapes = st["shapes"]
            effect = st.get("effect", "fade")
            gap = st.get("gap_ms", 180)
            for i, name in enumerate(shapes):
                report.append(f"slide {idx_str} [{name}]: {add_entrance(slide, name, effect, 'after_previous', i * gap)}")
        else:
            for a in anims:
                report.append(f"slide {idx_str} [{a.get('shape')}]: {add_entrance(slide, a['shape'], a.get('effect', 'fade'), a.get('trigger', 'after_previous'), a.get('delay_ms', 0))}")
    return report


def main():
    ap = argparse.ArgumentParser(description="Injeta entrance animations e transições em um .pptx")
    ap.add_argument("pptx")
    ap.add_argument("--entrance", nargs="+", metavar="ARG",
                    help="SHAPE EFFECT [TRIGGER] [DELAY_MS] — effect: appear|fade|zoom; trigger: on_click|after_previous|with_previous (default after_previous)")
    ap.add_argument("--transition", choices=sorted(TRANSITIONS))
    ap.add_argument("--slide", default=None, help="índice do slide para --transition (default: todos)")
    ap.add_argument("--spec", help="arquivo JSON de spec de motion")
    ap.add_argument("--out", required=True)
    args = ap.parse_args()

    prs = Presentation(args.pptx)
    report = []

    if args.spec:
        with open(args.spec) as f:
            spec = json.load(f)
        report.extend(apply_motion_spec(prs, spec))

    if args.entrance:
        shape, effect = args.entrance[0], args.entrance[1]
        trigger = args.entrance[2] if len(args.entrance) > 2 else "after_previous"
        delay_ms = int(args.entrance[3]) if len(args.entrance) > 3 else 0
        for slide in prs.slides:
            report.append(add_entrance(slide, shape, effect, trigger, delay_ms))

    if args.transition:
        if args.slide and args.slide != "all":
            idx = int(args.slide)
            report.append(f"slide {idx}: {add_transition(prs.slides[idx - 1], args.transition)}")
        else:
            for i, slide in enumerate(prs.slides, start=1):
                report.append(f"slide {i}: {add_transition(slide, args.transition)}")

    prs.save(args.out)
    print("\n".join(report))
    print(f"Salvo em: {args.out}")


if __name__ == "__main__":
    sys.exit(main())
