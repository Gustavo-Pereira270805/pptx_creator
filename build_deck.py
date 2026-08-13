"""Gera a apresentação corporativa 'Ocean Depths' (tema theme-factory)."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ---- Paleta do tema Ocean Depths ----
NAVY = RGBColor(0x1A, 0x23, 0x32)
TEAL = RGBColor(0x2D, 0x8B, 0x8B)
SEAFOAM = RGBColor(0xA8, 0xDA, 0xDC)
CREAM = RGBColor(0xF1, 0xFA, 0xEE)
CARD = RGBColor(0x22, 0x30, 0x3F)          # painel sutil sobre o navy
CARD_LINE = RGBColor(0x2E, 0x40, 0x54)     # borda discreta dos painéis

FONT = "DejaVu Sans"
FONT_BOLD = "DejaVu Sans"

SW, SH = Inches(13.333), Inches(7.5)

prs = Presentation()
prs.slide_width = SW
prs.slide_height = SH
BLANK = prs.slide_layouts[6]


def slide_bg(slide, color=NAVY):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def rect(slide, l, t, w, h, fill=None, line=None, shape=MSO_SHAPE.RECTANGLE,
         line_w=None, radius=None):
    sp = slide.shapes.add_shape(shape, l, t, w, h)
    if radius is not None:
        try:
            sp.adjustments[0] = radius
        except Exception:
            pass
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid()
        sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = line_w or Pt(1)
    sp.shadow.inherit = False
    return sp


def text(slide, l, t, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
         space_after=0, line_spacing=1.0):
    """runs: lista de parágrafos; cada parágrafo é uma lista de (texto, size, bold, color)."""
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after)
        p.line_spacing = line_spacing
        for (txt, size, bold, color) in para:
            r = p.add_run()
            r.text = txt
            r.font.size = Pt(size)
            r.font.bold = bold
            r.font.name = FONT
            r.font.color.rgb = color
    return tb


def header(slide, title, kicker=None, page=None):
    rect(slide, Inches(0.65), Inches(0.62), Inches(0.09), Inches(0.62), fill=TEAL)
    y = Inches(0.55)
    if kicker:
        text(slide, Inches(0.95), y, Inches(11.0), Inches(0.3),
             [[(kicker.upper(), 11, True, TEAL)]], space_after=0)
        y = Inches(0.92)
    text(slide, Inches(0.95), y, Inches(11.2), Inches(0.75),
         [[(title, 30, True, CREAM)]])
    # rodapé
    text(slide, Inches(0.65), Inches(7.05), Inches(6), Inches(0.3),
         [[("OCEAN DEPTHS", 9, True, TEAL), ("  ·  demonstração de design", 9, False, SEAFOAM)]])
    if page is not None:
        text(slide, Inches(12.0), Inches(7.05), Inches(0.7), Inches(0.3),
             [[(f"{page:02d}", 9, True, TEAL)]], align=PP_ALIGN.RIGHT)
    corner_waves(slide)


def card(slide, l, t, w, h):
    return rect(slide, l, t, w, h, fill=CARD, line=CARD_LINE,
                shape=MSO_SHAPE.ROUNDED_RECTANGLE, line_w=Pt(1), radius=0.055)


WAVE = RGBColor(0x23, 0x34, 0x4A)  # onda decorativa sutil sobre o navy


def corner_waves(slide):
    """Ondas discretas de marca no canto inferior direito (slides de conteúdo)."""
    y0 = Inches(6.62)
    for i, w in enumerate([0.85, 1.55, 2.45]):
        rect(slide, Inches(12.55 - w), y0 - Inches(i * 0.07), Inches(w),
             Emu(int(0.025 * 914400)), fill=WAVE)


def icon(slide, kind, l, t, d):
    """Ícones geométricos lineares no estilo do tema (contorno teal)."""
    d = Inches(d)
    LW = Pt(1.5)
    if kind == "globe":
        rect(slide, l, t, d, d, fill=None, line=TEAL, line_w=LW, shape=MSO_SHAPE.OVAL)
        rect(slide, l + Emu(int(d * 0.08)), t + Emu(int(d * 0.47)), Emu(int(d * 0.84)),
             Emu(int(d * 0.06)), fill=TEAL)
    elif kind == "ruler":
        rect(slide, l + Emu(int(d * 0.35)), t, Emu(int(d * 0.2)), d, fill=TEAL)
        for frac in (0.22, 0.5, 0.78):
            rect(slide, l, t + Emu(int(d * frac)), Emu(int(d * 0.42)),
                 Emu(int(d * 0.06)), fill=TEAL)
    elif kind == "drop":
        sp = rect(slide, l, t, d, d, fill=None, line=TEAL, line_w=LW, shape=MSO_SHAPE.TEAR)
        rect(slide, l + Emu(int(d * 0.35)), t + Emu(int(d * 0.55)), Emu(int(d * 0.3)),
             Emu(int(d * 0.06)), fill=TEAL)
    elif kind == "pressure":
        rect(slide, l, t, d, d, fill=None, line=TEAL, line_w=LW, shape=MSO_SHAPE.OVAL)
        c = Emu(int(d / 2))
        for frac in (0.12, 0.5, 0.88):
            rect(slide, l + Emu(int(d * frac)) - Emu(int(d * 0.03)), t - Emu(int(d * 0.12)),
                 Emu(int(d * 0.06)), Emu(int(d * 0.24)), fill=TEAL)
        rect(slide, l + c - Emu(int(d * 0.03)), t + Emu(int(d * 0.45)), Emu(int(d * 0.06)),
             Emu(int(d * 0.12)), fill=TEAL)
    elif kind == "thermo":
        rect(slide, l + Emu(int(d * 0.42)), t, Emu(int(d * 0.16)), Emu(int(d * 0.6)), fill=TEAL)
        rect(slide, l + Emu(int(d * 0.25)), t + Emu(int(d * 0.55)), Emu(int(d * 0.5)),
             Emu(int(d * 0.4)), fill=None, line=TEAL, line_w=LW, shape=MSO_SHAPE.OVAL)


# =====================================================================
# SLIDE 1 — CAPA
# =====================================================================
s = prs.slides.add_slide(BLANK)
slide_bg(s)

# barra vertical de acento à esquerda
rect(s, Inches(0), Inches(0), Inches(0.14), SH, fill=TEAL)

# camadas de profundidade decorativas (rodapé)
layers = [
    (Inches(6.35), 0.08, RGBColor(0x2D, 0x8B, 0x8B)),
    (Inches(6.15), 0.42, RGBColor(0x26, 0x6F, 0x7A)),
    (Inches(5.95), 0.73, RGBColor(0x1F, 0x4F, 0x63)),
    (Inches(5.75), 1.15, RGBColor(0x16, 0x30, 0x44)),
]
for x, w, c in layers:
    rect(s, x, Inches(7.5 - 0.18), w, Inches(0.18), fill=c)
text(s, Inches(11.15), Inches(6.6), Inches(2.0), Inches(0.7),
     [[("0 m", 10, False, SEAFOAM)], [("6.000 m", 10, False, SEAFOAM)]],
     align=PP_ALIGN.RIGHT, space_after=2)

# bloco de título
text(s, Inches(0.95), Inches(2.0), Inches(11.4), Inches(0.4),
     [[("PESQUISA & EXPLORAÇÃO  ·  RELATÓRIO 2026", 12, True, TEAL)]])
text(s, Inches(0.95), Inches(2.45), Inches(11.4), Inches(1.5),
     [[("Ocean Depths", 60, True, CREAM)]])
rect(s, Inches(1.0), Inches(3.75), Inches(2.6), Inches(0.035), fill=SEAFOAM)
text(s, Inches(0.95), Inches(4.05), Inches(9.6), Inches(1.2),
     [[("Explorando o último continente inexplorado da Terra — as profundezas "
        "do oceano concentram os maiores mistérios do planeta.", 19, False, SEAFOAM)]],
     line_spacing=1.25)
text(s, Inches(0.95), Inches(6.15), Inches(9), Inches(0.4),
     [[("OCEAN RESEARCH GROUP", 13, True, CREAM),
       ("   |   Departamento de Exploração Submarina", 13, False, SEAFOAM)]])

# =====================================================================
# SLIDE 2 — CONTEXTO: O MAIOR HABITAT DA TERRA
# =====================================================================
s = prs.slides.add_slide(BLANK)
slide_bg(s)
header(s, "O maior habitat da Terra", kicker="Contexto", page=2)

text(s, Inches(0.95), Inches(1.7), Inches(11.4), Inches(0.85),
     [[("O oceano domina o planeta — e continua sendo a região menos explorada "
        "da superfície terrestre.", 16, False, SEAFOAM)]], line_spacing=1.2)

stats = [
    ("71%", "da superfície da Terra\né coberta por oceano", 1.0, "globe"),
    ("3.700 m", "é a profundidade média\ndos oceanos do planeta", 4.0, "ruler"),
    ("97%", "de toda a água do planeta\nestá nos oceanos", 7.0, "drop"),
]
for big, small, x, ic in stats:
    card(s, Inches(x), Inches(2.75), Inches(3.55), Inches(2.5))
    icon(s, ic, Inches(x + 2.78), Inches(3.05), 0.48)
    text(s, Inches(x + 0.25), Inches(3.05), Inches(3.05), Inches(1.0),
         [[(big, 40, True, TEAL)]])
    rect(s, Inches(x + 0.28), Inches(4.05), Inches(0.6), Inches(0.03), fill=SEAFOAM)
    lines = small.split("\n")
    text(s, Inches(x + 0.25), Inches(4.25), Inches(3.05), Inches(0.9),
         [[(ln, 12.5, False, CREAM)] for ln in lines], space_after=2, line_spacing=1.1)

text(s, Inches(0.95), Inches(5.7), Inches(11.4), Inches(0.6),
     [[("Menos de 5% das profundezas abaixo de 200 m já foram exploradas de "
        "forma direta.", 15, True, SEAFOAM)]])

# =====================================================================
# SLIDE 3 — PERFIL DAS ZONAS DE PROFUNDIDADE
# =====================================================================
s = prs.slides.add_slide(BLANK)
slide_bg(s)
header(s, "O perfil vertical do oceano", kicker="Zonas de profundidade", page=3)

text(s, Inches(0.95), Inches(1.55), Inches(11.0), Inches(0.5),
     [[("A coluna de água é dividida em cinco zonas, cada uma com um regime "
        "próprio de luz, pressão e vida.", 13.5, False, SEAFOAM)]], line_spacing=1.15)

# área do "perfil" (faixas horizontais proporcionais à profundidade)
area_x, area_y, area_w, area_h = Inches(3.1), Inches(2.05), Inches(4.7), Inches(3.75)
scale = 3.75 / 11000.0  # polegadas por metro

zones = [
    ("Epipelágica", "0 – 200 m", 0, 200, RGBColor(0x4F, 0xB3, 0xB3)),
    ("Mesopelágica", "200 – 1.000 m", 200, 1000, RGBColor(0x2D, 0x8B, 0x8B)),
    ("Batipelágica", "1.000 – 4.000 m", 1000, 4000, RGBColor(0x22, 0x64, 0x73)),
    ("Abissopelágica", "4.000 – 6.000 m", 4000, 6000, RGBColor(0x1B, 0x45, 0x58)),
    ("Hadalpelágica", "6.000 – 10.935 m", 6000, 10935, RGBColor(0x14, 0x2C, 0x3E)),
]

# borda da área
rect(s, area_x, area_y, area_w, area_h, fill=None, line=TEAL, line_w=Pt(1.25))

y = area_y
for name, rng, d0, d1, color in zones:
    h = Emu(int((d1 - d0) * scale * 914400))
    rect(s, area_x, y, area_w, h, fill=color)
    # rótulos à esquerda da área (com folga para não colidir com marcadores)
    label_x = Inches(0.9)
    text(s, label_x, y + Emu(int(h / 2)) - Inches(0.14), Inches(1.95), Inches(0.55),
         [[(name, 11.5, True, CREAM)], [(rng, 10, False, SEAFOAM)]], space_after=0,
         align=PP_ALIGN.RIGHT)
    y += h

# marca da linha de luz (1.000 m) — contida dentro da área de plotagem
light_y = area_y + Emu(int(1000 * scale * 914400))
rect(s, area_x, light_y, area_w, Emu(int(0.02 * 914400)), fill=SEAFOAM)
text(s, Inches(7.95), light_y - Inches(0.16), Inches(4.5), Inches(0.35),
     [[("▲", 11, True, SEAFOAM), ("  limite da luz solar — 1.000 m", 11, True, SEAFOAM)]])

text(s, Inches(0.95), Inches(6.15), Inches(11.4), Inches(0.6),
     [[("Abaixo de 1.000 m, a escuridão é total: a fotossíntese não existe e "
        "a vida depende de neve marinha e quimiossíntese.", 13, False, CREAM)]],
     line_spacing=1.2)

# =====================================================================
# SLIDE 4 — NÚMEROS QUE IMPRESSIONAM
# =====================================================================
s = prs.slides.add_slide(BLANK)
slide_bg(s)
header(s, "Números que impressionam", kicker="Dados-chave", page=4)

data = [
    ("10.935 m", "profundidade da Fossa das\nMarianas (Challenger Deep)", 1.0, "ruler"),
    ("1.086 atm", "pressão no ponto mais\nprofundo do oceano", 4.55, "pressure"),
    ("2 – 4 °C", "temperatura média das\náguas profundas", 8.1, "thermo"),
]
for big, small, x, ic in data:
    card(s, Inches(x), Inches(1.85), Inches(4.2), Inches(2.15))
    icon(s, ic, Inches(x + 3.42), Inches(2.12), 0.55)
    text(s, Inches(x + 0.3), Inches(2.1), Inches(3.6), Inches(0.95),
         [[(big, 34, True, TEAL)]])
    lines = small.split("\n")
    text(s, Inches(x + 0.3), Inches(3.15), Inches(3.6), Inches(0.8),
         [[(ln, 12, False, CREAM)] for ln in lines], space_after=2, line_spacing=1.1)

# faixa de mapeamento
card(s, Inches(1.0), Inches(4.35), Inches(11.3), Inches(1.7))
rect(s, Inches(1.0), Inches(4.35), Inches(0.07), Inches(1.7), fill=TEAL)
text(s, Inches(1.35), Inches(4.6), Inches(10.6), Inches(0.6),
     [[("~25%", 26, True, TEAL),
       ("  do fundo do oceano está mapeado em alta resolução.", 14, False, CREAM)]])
text(s, Inches(1.35), Inches(5.35), Inches(10.6), Inches(0.6),
     [[("Aproximadamente 75% do assoalho oceânico ainda não tem mapeamento "
        "batimétrico de qualidade — menos do que a superfície de Marte.",
        12.5, False, SEAFOAM)]], line_spacing=1.2)

text(s, Inches(0.95), Inches(6.3), Inches(11.4), Inches(0.5),
     [[("O Monte Everest (8.849 m) caberia na Fossa das Marianas com "
        "mais de 2 km de sobra.", 13, True, SEAFOAM)]])

# =====================================================================
# SLIDE 5 — O AMBIENTE MAIS HOSTIL DA TERRA
# =====================================================================
s = prs.slides.add_slide(BLANK)
slide_bg(s)
header(s, "O ambiente mais hostil da Terra", kicker="Desafios físicos", page=5)

challenges = [
    ("P", "Pressão esmagadora", "A pressão aumenta 1 atm a cada 10 m. "
     "No Challenger Deep, equivale a 1.000 toneladas por metro quadrado."),
    ("L", "Escuridão absoluta", "A luz solar não atravessa 1.000 m. "
     "As criaturas profundas produzem a própria luz por bioluminescência."),
    ("T", "Frio constante", "Entre 2 e 4 °C, sem variação sazonal. "
     "A vida depende de quimiossíntese nas fontes hidrotermais."),
]
x = 1.0
for letter, title, body in challenges:
    card(s, Inches(x), Inches(1.95), Inches(3.65), Inches(4.3))
    circ = rect(s, Inches(x + 0.28), Inches(2.3), Inches(0.85), Inches(0.85),
                fill=TEAL, shape=MSO_SHAPE.OVAL)
    text(s, Inches(x + 0.28), Inches(2.42), Inches(0.85), Inches(0.6),
         [[(letter, 24, True, NAVY)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text(s, Inches(x + 0.28), Inches(3.4), Inches(3.1), Inches(0.8),
         [[(title, 16, True, CREAM)]], line_spacing=1.1)
    rect(s, Inches(x + 0.3), Inches(3.95), Inches(0.55), Inches(0.03), fill=SEAFOAM)
    text(s, Inches(x + 0.28), Inches(4.15), Inches(3.1), Inches(1.9),
         [[(body, 12, False, SEAFOAM)]], line_spacing=1.25)
    x += 3.85

text(s, Inches(0.95), Inches(6.5), Inches(11.4), Inches(0.5),
     [[("Sobreviver aqui exige engenharia de precisão — e os mesmos desafios "
        "impulsionam a inovação.", 13, True, CREAM)]])

# =====================================================================
# SLIDE 6 — TECNOLOGIA DE EXPLORAÇÃO
# =====================================================================
s = prs.slides.add_slide(BLANK)
slide_bg(s)
header(s, "Tecnologia de exploração", kicker="Como chegamos lá", page=6)

techs = [
    ("Sonar multifeixe", "Mapeia o fundo do oceano emitindo feixes acústicos. "
     "É a base da cartografia batimétrica moderna."),
    ("ROVs", "Robôs operados à distância alcançam mais de 6.000 m, "
     "coletando imagens e amostras sem risco humano."),
    ("Submersíveis", "Naves tripuladas como o Limiting Factor (10.925 m) "
     "levaram humanos ao ponto mais profundo em 2019."),
]
y = 1.85
for i, (title, body) in enumerate(techs, start=1):
    card(s, Inches(1.0), Inches(y), Inches(11.3), Inches(1.35))
    rect(s, Inches(1.0), Inches(y), Inches(0.07), Inches(1.35), fill=TEAL)
    text(s, Inches(1.35), Inches(y + 0.18), Inches(10.5), Inches(0.4),
         [[(f"0{i}", 13, True, TEAL), ("   ", 13, False, CREAM), (title, 17, True, CREAM)]])
    text(s, Inches(1.35), Inches(y + 0.62), Inches(10.4), Inches(0.65),
         [[(body, 12.5, False, SEAFOAM)]], line_spacing=1.2)
    y += 1.55

text(s, Inches(0.95), Inches(6.55), Inches(11.4), Inches(0.5),
     [[("1960 — o batiscafo Trieste foi o primeiro a tocar o Challenger Deep.",
        12.5, False, CREAM)]])

# =====================================================================
# SLIDE 7 — MAPEAR O DESCONHECIDO
# =====================================================================
s = prs.slides.add_slide(BLANK)
slide_bg(s)
header(s, "Mapear o desconhecido", kicker="A próxima fronteira", page=7)

card(s, Inches(1.0), Inches(1.85), Inches(5.4), Inches(4.2))
text(s, Inches(1.3), Inches(2.35), Inches(4.8), Inches(1.6),
     [[("75%", 66, True, TEAL)]])
text(s, Inches(1.3), Inches(4.05), Inches(4.8), Inches(1.0),
     [[("do assoalho oceânico ainda não foi mapeado em alta resolução — "
        "menos conhecido do que a superfície de Marte.", 14, False, CREAM)]],
     line_spacing=1.25)
rect(s, Inches(1.32), Inches(5.35), Inches(0.6), Inches(0.03), fill=SEAFOAM)

points = [
    ("Projeto Seabed 2030", "Iniciativa global com o objetivo de mapear 100% "
     "do fundo do oceano até 2030."),
    ("Por que importa", "Cabos submarinos, clima, rotas de navegação e "
     "recursos minerais dependem de dados batimétricos."),
    ("Alto retorno", "Um oceano mapeado reduz riscos e abre caminho para "
     "conservação baseada em dados."),
]
y = 1.85
for title, body in points:
    card(s, Inches(6.7), Inches(y), Inches(5.6), Inches(1.25))
    text(s, Inches(6.95), Inches(y + 0.14), Inches(5.1), Inches(0.35),
         [[(title.upper(), 11.5, True, TEAL)]])
    text(s, Inches(6.95), Inches(y + 0.52), Inches(5.1), Inches(0.65),
         [[(body, 11.5, False, CREAM)]], line_spacing=1.15)
    y += 1.42

# =====================================================================
# SLIDE 8 — ENCERRAMENTO
# =====================================================================
s = prs.slides.add_slide(BLANK)
slide_bg(s)
rect(s, Inches(0), Inches(0), Inches(0.14), SH, fill=TEAL)

# camadas de profundidade no rodapé
for x, w, c in layers:
    rect(s, x, Inches(7.5 - 0.18), w, Inches(0.18), fill=c)

text(s, Inches(0.95), Inches(2.1), Inches(11.4), Inches(0.4),
     [[("OCEAN RESEARCH GROUP", 12, True, TEAL)]])
text(s, Inches(0.95), Inches(2.55), Inches(11.4), Inches(1.3),
     [[("O oceano espera.", 54, True, CREAM)]])
rect(s, Inches(1.0), Inches(3.85), Inches(2.6), Inches(0.035), fill=SEAFOAM)
text(s, Inches(0.95), Inches(4.15), Inches(10.4), Inches(1.0),
     [[("A próxima fronteira da exploração começa abaixo de 1.000 metros — "
        "e as ferramentas para alcançá-la nunca estiveram tão maduras.",
        18, False, SEAFOAM)]], line_spacing=1.25)
text(s, Inches(0.95), Inches(6.15), Inches(9), Inches(0.4),
     [[("ocean-research@corp.com", 13, False, CREAM),
       ("    |    www.oceandepths.corp", 13, False, SEAFOAM)]])

prs.save("ocean-depths.pptx")
print("Deck salvo: ocean-depths.pptx com", len(prs.slides.__iter__.__self__._sldIdLst), "slides")
