#!/usr/bin/env python3
"""Builds the redesigned Capacitação SIME deck (17 slides) using the
figma-pitch-deck design system (design/design-spec.md). Content is EXACTLY
the original (texts/numbers/data preserved); only the visual chrome changed.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ---------------------------------------------------------------- tokens
PAPER = "F2F3F2"
WHITE = "FFFFFF"
NAVY = "0E2756"
SLATE = "304258"
BLUE = "2353E9"
BRIGHT = "559AFF"
ROYAL = "193BA7"
STEEL = "6E7885"
TINT = "EAF1FF"
AMBER = "B4540A"
FONT = "Poppins"

MEDIA = "media_original"
OUT = "Capacitacao_SIME_Sebrae_PB_redesign.pptx"

prs = Presentation()
prs.slide_width = Emu(12192000)   # 13.333 in
prs.slide_height = Emu(6858000)   # 7.5 in
BLANK = prs.slide_layouts[6]


def C(h): return RGBColor.from_string(h)


def _set_font(run, size, bold, color, name=FONT):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = C(color)
    run.font.name = name


def _alpha(shape, opacity_pct):
    sF = shape.fill._xPr.find(qn('a:solidFill'))
    clr = sF.find(qn('a:srgbClr'))
    a = clr.makeelement(qn('a:alpha'), {'val': str(int(opacity_pct * 1000))})
    clr.append(a)


def new_slide():
    s = prs.slides.add_slide(BLANK)
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    r.fill.solid(); r.fill.fore_color.rgb = C(PAPER)
    r.line.fill.background(); r.shadow.inherit = False
    return s


def add_rect(slide, x, y, w, h, fill=WHITE, line=None, lw=1.0, radius=None):
    kind = MSO_SHAPE.ROUNDED_RECTANGLE if radius is not None else MSO_SHAPE.RECTANGLE
    sp = slide.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
    if radius is not None:
        try:
            sp.adjustments[0] = radius
        except Exception:
            pass
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid(); sp.fill.fore_color.rgb = C(fill)
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = C(line); sp.line.width = Pt(lw)
    sp.shadow.inherit = False
    return sp


def add_text(slide, x, y, w, h, text, size=11, bold=False, color=STEEL,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, spacing=None, after=None, name=None):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    if name:
        tb.name = name
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    if isinstance(text, str):
        text = [[(text, {})]]
    elif isinstance(text, list) and text and isinstance(text[0], tuple):
        text = [text]  # single paragraph as list of (run_text, opts)
    elif isinstance(text, list) and text and isinstance(text[0], str):
        text = [[(t, {}) for t in text]]  # multiple plain paragraphs
    first = True
    for para in text:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = align
        if spacing: p.line_spacing = spacing
        if after: p.space_after = Pt(after)
        for run_text, opts in para:
            r = p.add_run(); r.text = run_text
            _set_font(r, opts.get('size', size), opts.get('bold', bold),
                      opts.get('color', color), opts.get('font', FONT))
    return tb


def add_pic_cover(slide, img, x, y, w, h, name=None):
    pic = slide.shapes.add_picture(img, Inches(x), Inches(y), Inches(w), Inches(h))
    if name:
        pic.name = name
    try:
        iw, ih = pic.image.size
        fr, ir = w / h, iw / ih
        if ir > fr:
            c = (1 - fr / ir) / 2
            pic.crop_left = c; pic.crop_right = c
        else:
            c = (1 - ir / fr) / 2
            pic.crop_top = c; pic.crop_bottom = c
    except Exception:
        pass
    pic.line.fill.background()
    return pic


def accent_square(slide, x, y, s=0.14, color=BLUE):
    add_rect(slide, x, y, s, s, fill=NAVY, radius=None)
    add_rect(slide, x + 0.05, y + 0.05, s, s, fill=color, radius=None)


def eyebrow(slide, x, y, text, color=BLUE, size=10.5):
    accent_square(slide, x, y + 0.02, 0.13, color)
    add_text(slide, x + 0.24, y - 0.02, 8.0, 0.3, text.upper(), size=size, bold=True, color=color)


def header(slide, eyebrow_text, title, meta="Sebrae/PB", title_size=24):
    eyebrow(slide, 0.5, 0.42, eyebrow_text)
    add_text(slide, 0.5, 0.74, 11.5, 0.6, title, size=title_size, bold=True, color=NAVY, name="title")
    add_text(slide, 11.2, 0.5, 1.63, 0.3, meta, size=10, bold=True, color=ROYAL, align=PP_ALIGN.RIGHT, name="meta")


def footer(slide, meta="Monitoramento da Estratégia • Sebrae/PB", num=None, total=17):
    add_rect(slide, 0.5, 7.16, 12.33, 0.014, fill="D5DAE3")
    add_text(slide, 0.5, 7.23, 8.0, 0.22, meta, size=8, color=STEEL)
    if num:
        add_text(slide, 11.6, 7.23, 1.23, 0.22, f"Slide {num} / {total}", size=8,
                 color=STEEL, align=PP_ALIGN.RIGHT)


def card(slide, x, y, w, h, fill=WHITE, line="E3E7EE", radius=0.05):
    return add_rect(slide, x, y, w, h, fill=fill, line=line, lw=1.0, radius=radius)


def chip(slide, x, y, w, h, text, fill=BLUE, color=WHITE, size=11, line=None):
    add_rect(slide, x, y, w, h, fill=fill, line=line, lw=1.0, radius=0.5)
    add_text(slide, x, y + (h - 0.22) / 2, w, 0.26, text, size=size, bold=True,
             color=color, align=PP_ALIGN.CENTER)


def warning(slide, x, y, w, text, h=0.72, prefix="⚠️ Atenção (Erro Comum): "):
    card(slide, x, y, w, h, fill=TINT, line=None, radius=0.06)
    add_rect(slide, x, y + 0.08, 0.05, h - 0.16, fill=BLUE)
    add_text(slide, x + 0.22, y + 0.1, w - 0.4, h - 0.2,
             [[(prefix, {'bold': True, 'color': NAVY}),
               (text, {'color': AMBER})]],
             size=9.5, spacing=1.08)


def tip_bar(slide, x, y, w, prefix, text, h=0.5):
    card(slide, x, y, w, h, fill=TINT, line=None, radius=0.07)
    add_rect(slide, x, y + 0.07, 0.05, h - 0.14, fill=BLUE)
    add_text(slide, x + 0.22, y + (h - 0.3) / 2, w - 0.44, h - 0.1,
             [[(prefix + " ", {'bold': True, 'color': NAVY}), (text, {'color': SLATE})]],
             size=9.5, anchor=MSO_ANCHOR.MIDDLE)


def quote_bar(slide, x, y, w, quote, tag=None, h=0.62, qsize=13.5):
    card(slide, x, y, w, h, fill=NAVY, line=None, radius=0.06)
    add_rect(slide, x + 0.22, y + 0.12, 0.05, h - 0.24, fill=BRIGHT)
    runs = [[("“", {'bold': True, 'color': BRIGHT, 'size': qsize + 4}),
             (quote, {'bold': False, 'color': WHITE, 'size': qsize})]]
    if tag:
        runs.append([(tag, {'bold': True, 'color': BRIGHT, 'size': 10})])
    add_text(slide, x + 0.42, y + 0.12, w - 0.75, h - 0.2, runs, spacing=1.1,
             anchor=MSO_ANCHOR.MIDDLE)


def table_slide(slide, x, y, w, rows_data, col_w, header_bg=NAVY, header_h=0.42,
                row_h=0.62, font_size=9.5, header_size=9.5):
    n_rows = len(rows_data)
    n_cols = len(rows_data[0])
    total_h = header_h + row_h * (n_rows - 1)
    gfx = slide.shapes.add_table(n_rows, n_cols, Inches(x), Inches(y), Inches(w), Inches(total_h))
    tbl = gfx.table
    tbl.first_row = False; tbl.horz_banding = False
    for ci, cw in enumerate(col_w):
        tbl.columns[ci].width = Inches(cw)
    for ri, row in enumerate(rows_data):
        tbl.rows[ri].height = Inches(header_h if ri == 0 else row_h)
        for ci, val in enumerate(row):
            cell = tbl.cell(ri, ci)
            cell.margin_left = Inches(0.08); cell.margin_right = Inches(0.06)
            cell.margin_top = Inches(0.04); cell.margin_bottom = Inches(0.04)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.fill.solid(); cell.fill.fore_color.rgb = C(header_bg if ri == 0 else (WHITE if ri % 2 == 1 else "F7F8FA"))
            tf = cell.text_frame; tf.word_wrap = True
            p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
            r = p.add_run(); r.text = val
            _set_font(r, header_size if ri == 0 else font_size,
                      True, WHITE if ri == 0 else (NAVY if ci == 0 else SLATE))
    return tbl, gfx


def kpi_card(slide, x, y, w, h, eyebrow_t, label, number, note, img_path=None, num_size=60, name=None):
    card(slide, x, y, w, h, fill=WHITE, line="E3E7EE", radius=0.05)
    if name:
        slide.shapes[-1].name = name
    add_rect(slide, x, y, w, 0.09, fill=BLUE)
    add_text(slide, x + 0.3, y + 0.24, w - 0.6, 0.26, eyebrow_t.upper(), size=9.5, bold=True, color=BLUE, name="kpi_eyebrow")
    add_text(slide, x + 0.3, y + 0.55, w - 0.6, 0.6, label, size=13, bold=True, color=NAVY, spacing=1.05, name="kpi_label")
    ph = number.startswith("[")
    add_text(slide, x + 0.3, y + 1.35, w - 0.6, 0.95, number, size=(22 if ph else num_size),
             bold=True, color=(STEEL if ph else BLUE), name="kpi_number")
    add_text(slide, x + 0.3, y + 2.45, w - 0.6, 0.7, note, size=8.5, color=STEEL, spacing=1.05, name="kpi_note")
    if img_path:
        add_pic_cover(slide, img_path, x + 0.25, y + 3.15, w - 0.5, h - 3.35, name="kpi_img")


def meta_page(idx, eyebrow_t, title, blocks, attention, kpi, img_left, img_right, title_size=22):
    s = new_slide()
    header(s, f"Meta Mobilizadora {idx:02d} / 06", title, title_size=title_size)
    lx, lw = 0.5, 8.0
    y = 1.62
    for bi, (label, body) in enumerate(blocks, 1):
        add_text(s, lx, y, lw, 0.28, label, size=11.5, bold=True, color=NAVY, name=f"q{bi}")
        add_text(s, lx, y + 0.32, lw, 0.5, body, size=9.8, color=SLATE, spacing=1.1, name=f"a{bi}")
        y += 0.32 + 0.5 + 0.14
    if attention:
        warning(s, lx, y, lw, attention, h=0.8)
    add_pic_cover(s, img_left, lx, 6.34, lw, 0.7, name="img_left")
    kpi_card(s, 8.72, 1.62, 4.11, 5.5, "Meta Sebrae/PB 2026", *kpi, img_path=img_right, name="kpi_card")
    return s


# ================================================================ SLIDE 1 — Capa
s = new_slide()
s.shapes.add_picture(f"{MEDIA}/image1.jpeg", 0, 0, prs.slide_width, prs.slide_height)
ov = add_rect(s, 0, 0, 9.0, 7.5, fill=NAVY, radius=None)
_alpha(ov, 82)
add_rect(s, 0.5, 1.05, 0.5, 0.055, fill=BRIGHT)
add_text(s, 0.5, 1.35, 8.5, 0.3, "Sistema de Monitoramento da Estratégia", size=12, bold=True, color=BRIGHT)
add_text(s, 0.5, 2.0, 9.0, 1.0, "SIME", size=54, bold=True, color=WHITE)
add_text(s, 0.5, 2.95, 9.0, 0.8, "Entendendo nossos resultados", size=34, bold=True, color=WHITE)
add_text(s, 0.5, 4.15, 8.3, 0.55, "Capacitação interna Sebrae/PB sobre Painéis de Resultados de Atendimento.",
         size=13, color="C9D4E8", spacing=1.15)
add_rect(s, 0.5, 5.05, 3.1, 0.02, fill="3B4E77")
add_text(s, 0.5, 5.25, 8.0, 0.3, "Objetivo da nossa conversa", size=11.5, bold=True, color=BRIGHT)
add_text(s, 0.5, 5.62, 8.3, 1.1,
         "Entender o que cada indicador mede, de onde vêm os números e como interpretar corretamente cada painel de resultados estratégico.",
         size=12, color="E3E9F5", spacing=1.2)
add_text(s, 0.5, 7.1, 5.0, 0.22, "Sebrae Paraíba • Capacitação SIME", size=9, color="9FB0CC")

# ================================================================ SLIDE 2 — Conceitos
s = new_slide()
header(s, "Conceitos Fundamentais", "Antes de começar: Atendimento ≠ Cliente Atendido")
card(s, 0.5, 1.62, 3.94, 5.35)
add_text(s, 0.75, 1.85, 3.4, 0.26, "QUEM É CONTADO", size=9.5, bold=True, color=BLUE)
add_text(s, 0.75, 2.16, 3.4, 0.4, "Cliente atendido", size=15, bold=True, color=NAVY)
add_pic_cover(s, f"{MEDIA}/image2.png", 0.75, 2.75, 3.44, 2.6)
add_text(s, 0.75, 5.5, 3.44, 1.4,
         [[("Pessoa Física ou Pequeno Negócio (CNPJ ou CPF) que recebeu pelo menos um atendimento aceito dentro do período de referência analisado.", {})],
          [("A contagem é por cliente único.", {'bold': True, 'color': NAVY})]],
         size=9.3, color=SLATE, spacing=1.1, after=6)
card(s, 4.62, 1.62, 3.94, 5.35)
add_text(s, 4.87, 1.85, 3.4, 0.26, "O QUE É CONTADO", size=9.5, bold=True, color=BLUE)
add_text(s, 4.87, 2.16, 3.4, 0.4, "Atendimento", size=15, bold=True, color=NAVY)
add_pic_cover(s, f"{MEDIA}/image3.png", 4.87, 2.75, 3.44, 2.6)
add_text(s, 4.87, 5.5, 3.44, 1.1,
         [[("Cada interação realizada com o cliente. Representa o volume total de serviços e soluções prestadas no período.", {})],
          [("A contagem é por interação.", {'bold': True, 'color': NAVY})]],
         size=9.3, color=SLATE, spacing=1.1, after=6)
add_text(s, 4.87, 6.62, 3.44, 0.3, [("⚠️ Observação: ", {'bold': True, 'color': AMBER}),
                                    ("O instrumento Informação não gera atendimento.", {'color': SLATE})],
         size=8.5, spacing=1.05)
card(s, 8.74, 1.62, 4.09, 5.35)
add_text(s, 9.0, 1.86, 3.6, 0.26, "EXEMPLO PRÁTICO", size=9.5, bold=True, color=BLUE)
add_text(s, 9.0, 2.18, 3.6, 0.35, "Empresa A", size=13.5, bold=True, color=NAVY)
yy = 2.66
for it in ["1 Sebratec", "1 Empretec", "1 NF no portal Sebrae", "1 Palestra na FE"]:
    add_text(s, 9.0, yy, 3.5, 0.28, f"➔ {it}", size=10.5, color=SLATE)
    yy += 0.34
add_rect(s, 9.0, 4.3, 3.58, 0.02, fill="D5DAE3")
add_text(s, 9.0, 4.5, 3.5, 0.26, "Resultado consolidado", size=9.5, bold=True, color=BLUE)
card(s, 9.0, 4.86, 3.58, 0.95, fill=TINT, line=None, radius=0.08)
add_text(s, 9.25, 5.0, 3.1, 0.3, "1 Cliente Atendido", size=12.5, bold=True, color=NAVY)
add_text(s, 9.25, 5.32, 3.1, 0.3, "4 Atendimentos", size=12.5, bold=True, color=BLUE)
footer(s, num=2)

# ================================================================ SLIDE 3 — Painéis
s = new_slide()
header(s, "Regras de Contabilização", "Painéis: Clientes Atendidos e Atendimentos x Metas Mobilizadoras", title_size=20)
card(s, 0.5, 1.62, 5.7, 5.35)
add_text(s, 0.8, 1.86, 5.0, 0.26, "FOCO NA OPERAÇÃO", size=9.5, bold=True, color=BLUE)
add_text(s, 0.8, 2.18, 5.1, 0.35, "Painel de Clientes Atendidos e Atendimentos", size=13.5, bold=True, color=NAVY)
add_pic_cover(s, f"{MEDIA}/image3.png", 0.8, 2.7, 5.1, 2.0)
yy = 4.92
for b in ["Olha diretamente para a atuação operacional do Sebrae/PB.",
          "Inclui clientes de outras UFs que foram atendidos pelo Sebrae/PB.",
          "Inclui empresas paraibanas atendidas diretamente pelo Sebrae Nacional (Sebrae/NA)."]:
    add_text(s, 0.95, yy, 4.95, 0.5, [[("▪ ", {'bold': True, 'color': BLUE}), (b, {})]],
             size=9.3, color=SLATE, spacing=1.08)
    yy += 0.62
card(s, 7.13, 1.62, 5.7, 5.35)
add_text(s, 7.43, 1.86, 5.0, 0.26, "FOCO NO TERRITÓRIO", size=9.5, bold=True, color=BLUE)
add_text(s, 7.43, 2.18, 5.1, 0.35, "Painel de Metas Mobilizadoras", size=13.5, bold=True, color=NAVY)
add_pic_cover(s, f"{MEDIA}/image3.png", 7.43, 2.7, 5.1, 2.0)
yy = 4.92
for b in ["Olha com foco no público-alvo paraibano.",
          "Não contabiliza empresas com cadastro de outros estados, mesmo atendidas pelo PB.",
          "Contabiliza empresa paraibana que tenha sido atendida por outro Sebrae/UF no Brasil."]:
    add_text(s, 7.58, yy, 4.95, 0.5, [[("▪ ", {'bold': True, 'color': BLUE}), (b, {})]],
             size=9.3, color=SLATE, spacing=1.08)
    yy += 0.62
add_rect(s, 6.32, 3.9, 0.71, 0.71, fill=BLUE, line=None, radius=1.0)
add_text(s, 6.32, 4.08, 0.71, 0.4, "VS", size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
footer(s, num=3)

# ================================================================ SLIDE 4 — Matriz
s = new_slide()
header(s, "Matriz de Contabilização", "Um exemplo para não esquecer")
rows = [
    ["Caso / Cenário Prático", "Painel de Clientes Atendidos", "Painel de Metas Mobilizadoras"],
    ["Empresa da Paraíba atendida pelo Sebrae/PB", "SIM", "SIM"],
    ["Empresa de Pernambuco atendida pelo Sebrae/PB", "SIM", "NÃO"],
    ["Empresa da Paraíba atendida pelo Sebrae/PE", "NÃO", "SIM"],
    ["Empresa da Paraíba atendida pelo Sebrae Nacional", "SIM", "SIM"],
]
table_slide(s, 0.5, 1.7, 12.33, rows, col_w=[5.53, 3.4, 3.4], row_h=0.78, font_size=10.5, header_size=11)
for ci in (1, 2):
    x = 0.5 + 5.53 + (ci - 1) * 3.4
    cy = 1.7 + 0.42
    for ri in range(1, 5):
        val = rows[ri][ci]
        if val == "SIM":
            chip(s, x + (3.4 - 1.1) / 2, cy + (0.78 - 0.4) / 2, 1.1, 0.4, "✔ SIM", fill=BLUE, size=10.5)
        else:
            chip(s, x + (3.4 - 1.1) / 2, cy + (0.78 - 0.4) / 2, 1.1, 0.4, "✘ NÃO", fill="FFFFFF", line="C3CAD6", color=STEEL, size=10.5)
        cy += 0.78
footer(s, num=4)

# ================================================================ SLIDE 5 — As 6 Metas
s = new_slide()
header(s, "Direcionamento Estratégico", "As 6 Metas Mobilizadoras do Sebrae/PB")
metas = [
    ("Meta 1", "NPS", "Grau de recomendação do Sebrae pelo cliente", "image5.png"),
    ("Meta 2", "ME + EPP", "Percentual de Microempresas (ME) e Empresas de Pequeno Porte (EPP) atendidas pelo Sebrae no ano.", "image5.png"),
    ("Meta 3", "MEI", "Percentual de Microempreendedores Individuais (MEI) atendidos pelo Sebrae no ano.", "image5.png"),
    ("Meta 4", "Engajamento", "Percentual dos pequenos negócios atendidos pelo Sebrae considerados como engajados.", "image5.png"),
    ("Meta 5", "Grupos Sub-representados", "Proporção de atendimentos do Sebrae às pessoas dos grupos sub-representados em relação à proporção populacional por tipo de grupo sub-representado.", "image5.png"),
    ("Meta 6", "REB + Convênio", "Resultados gerados via Rede Estadual de Atendimento e parcerias em convênios de cooperação estratégica.", "image5.png"),
]
cw, ch, gap = 3.94, 2.55, 0.255
for i, (ml, title, desc, img) in enumerate(metas):
    col, row = i % 3, i // 3
    x = 0.5 + col * (cw + gap)
    y = 1.62 + row * (ch + 0.3)
    card(s, x, y, cw, ch)
    add_pic_cover(s, f"{MEDIA}/{img}", x, y, cw, 1.15)
    add_text(s, x + 0.28, y + 1.26, cw - 0.56, 0.24, ml.upper(), size=9, bold=True, color=BLUE)
    add_text(s, x + 0.28, y + 1.52, cw - 0.56, 0.32, title, size=13.5, bold=True, color=NAVY)
    add_text(s, x + 0.28, y + 1.86, cw - 0.56, 0.66, desc, size=8.4, color=SLATE, spacing=1.05)
footer(s, num=5)

# ================================================================ SLIDES 6-11 — Metas
meta_page(1, "Meta Mobilizadora 01 / 06", "NPS (Net Promoter Score)",
    [("O que queremos medir?", "Grau de recomendação do Sebrae pelo cliente"),
     ("Como é coletado?", "Pesquisa realizada de forma online, com envio de convite ao cliente por WhatsApp 60 dia após o registro do atendimento, tendo sido ampliada, a partir de 2024, para envios também via WhatsApp."),
     ("Como é calculado?", "% Clientes Promotores (notas 9-10) − % Clientes Detratores (notas 0-6)")],
    None,
    ("Índice NPS Alvo", "91,5", "* Valor parametrizado conforme acordado no planejamento estratégico anual."),
    f"{MEDIA}/image6.png", f"{MEDIA}/image7.png")
s = prs.slides[-1]
warning(s, 0.5, 5.42, 8.0, "Não são pesquisados clientes cujo atendimento é realizado pelos canais email e redes sociais.",
        h=0.58, prefix="⚠️ Observação: ")
footer(s, num=6)

meta_page(2, "Meta Mobilizadora 02 / 06", "ME + EPP (Microempresas e Empresas de Pequeno Porte)",
    [("O que queremos medir?", "A penetração e o alcance do Sebrae junto ao universo de Microempresas (ME) e Empresas de Pequeno Porte (EPP) do estado."),
     ("Quem entra na conta?", "ME e EPP sediadas no estado da Paraíba que foram atendidas por qualquer canal e unidade do Sistema Sebrae."),
     ("Como é calculado?", "Total de ME e EPP Atendidas (CNPJ Único) ÷ Universo Total de ME/EPP do Estado"),
     ("Desdobramento para as Regionais", "Distribuição de metas operacionais proporcional ao universo empresarial (densidade de empresas) de cada um dos territórios.")],
    "Contar o volume absoluto de atendimentos (interações) em vez de clientes individuais (CNPJs únicos). Se uma mesma ME/EPP realiza 10 consultorias, ela conta apenas uma vez para o alcance da meta de cobertura.",
    ("Percentual de Cobertura", "[ PREENCHER % OFICIAL ]", "* Indicador estratégico nacional com pactuação local junto ao Conselho Deliberativo Estadual."),
    f"{MEDIA}/image6.png", f"{MEDIA}/image7.png")
s = prs.slides[-1]
footer(s, num=7)

meta_page(3, "Meta Mobilizadora 03 / 06", "MEI (Microempreendedor Individual)",
    [("O que queremos medir?", "A cobertura e o alcance do Sebrae junto aos Microempreendedores Individuais (MEI) da Paraíba."),
     ("Quem entra na conta?", "MEIs ativos no estado da Paraíba que foram atendidos por qualquer canal do Sistema Sebrae."),
     ("Como é calculado?", "MEIs Atendidos (CNPJ Único) ÷ Total de MEIs Ativos no Estado"),
     ("Desdobramento para as Regionais", "Proporcional ao número total de MEIs ativos registrados em cada município pertencente ao território da regional.")],
    "Incluir MEIs de outras Unidades Federativas (UFs) atendidos fisicamente em agências locais da Paraíba. Esse painel valida estritamente a localização jurídica do CNPJ ativo na Paraíba.",
    ("Percentual de Cobertura MEI", "[ PREENCHER % OFICIAL ]", "* Meta de extrema relevância para a inclusão social e formalização de novos negócios."),
    f"{MEDIA}/image6.png", f"{MEDIA}/image7.png")
s = prs.slides[-1]
footer(s, num=8)

meta_page(4, "Meta Mobilizadora 04 / 06", "Engajamento",
    [("O que queremos medir?", "A intensidade e a profundidade do relacionamento dos clientes com o Sebrae (fidelização e jornada de soluções)."),
     ("Quem entra na conta?", "Clientes com múltiplas interações qualificadas e complementares no período (ex: realizar curso + consultoria)."),
     ("Como é calculado?", "Regra nacional de engajamento baseada em uma pontuação ponderada de acordo com o tipo e complexidade de cada interação realizada."),
     ("Desdobramento para as Regionais", "Acompanhamento sistemático e detalhado por carteira de clientes ativos e metas de retenção atribuídas a cada unidade.")],
    "Confundir o indicador de engajamento com a quantidade bruta de atendimentos dispersos. O engajamento exige a complementaridade e evolução do mesmo cliente ao longo de múltiplos canais.",
    ("Pontuação de Engajamento", "[ PREENCHER META OFICIAL ]", "* Métrica orientada à migração de atendimentos eventuais para trilhas integradas de soluções."),
    f"{MEDIA}/image6.png", f"{MEDIA}/image7.png")
s = prs.slides[-1]
footer(s, num=9)

meta_page(5, "Meta Mobilizadora 05 / 06", "Grupos Sub-representados",
    [("O que queremos medir?", "O alcance e a cobertura de atendimento do Sebrae junto a públicos e grupos prioritários definidos nacionalmente."),
     ("Quem entra na conta?", "Empresas e empreendedores que autodeclaram pertencimento aos grupos mapeados como prioritários ou sub-representados no ato do cadastro."),
     ("Como é calculado?", "Clientes Atendidos Pertencentes aos Grupos (CNPJ Único) ÷ Universo Mapeado do Grupo"),
     ("Desdobramento para as Regionais", "Distribuição com base nas características geodemográficas e na concentração territorial de cada grupo prioritário no estado.")],
    "Classificar clientes de forma equivocada sem a devida autodeclaração formal ou sem documentos comprobatórios exigidos nas regras técnicas de cadastramento do indicador.",
    ("Cobertura de Sub-representados", "[ PREENCHER % OFICIAL ]", "* Métrica vital para o cumprimento da função social e do compromisso de diversidade do Sebrae."),
    f"{MEDIA}/image6.png", f"{MEDIA}/image7.png")
s = prs.slides[-1]
footer(s, num=10)

meta_page(6, "Meta Mobilizadora 06 / 06", "REB + Convênio",
    [("O que queremos medir?", "Os resultados decorrentes da atuação da Rede Estadual de Atendimento (REB) em parceria com convênios estratégicos vigentes."),
     ("Quem entra na conta?", "Clientes atendidos e registrados via pontos e salas da REB e ações vinculadas formalmente a convênios e parcerias técnicas locais."),
     ("Como é calculado?", "[ Regra definida pela Diretoria Técnica — Aguardando Parametrização ]"),
     ("Desdobramento para as Regionais", "Divisão proporcional conforme a capilaridade, participação ativa das salas locais e representatividade das unidades no ecossistema regional.")],
    "Misturar os dados e atendimentos diretos das salas da REB com os atendimentos gerais da equipe interna das agências Sebrae sem a correta codificação de parceiros.",
    ("Indicador REB", "[ AGUARDANDO REGRAS ]", "* Indicador pactuado localmente de acordo com o plano operacional das Unidades de Atendimento."),
    f"{MEDIA}/image6.png", f"{MEDIA}/image7.png")
s = prs.slides[-1]
footer(s, num=11)

# ================================================================ SLIDE 12 — Navegação
s = new_slide()
s.shapes.add_picture(f"{MEDIA}/image8.jpeg", 0, 0, prs.slide_width, prs.slide_height)
ov = add_rect(s, 0, 0, 13.333, 7.5, fill=NAVY, radius=None)
_alpha(ov, 86)
add_text(s, 0.7, 0.62, 8.0, 0.3, "Navegação Prática", size=12, bold=True, color=BRIGHT)
add_text(s, 0.7, 1.5, 11.0, 0.9, "Agora vamos para o SIME!", size=40, bold=True, color=WHITE)
quote_bar(s, 0.7, 2.75, 9.6, "Agora que sabemos exatamente o que os números significam… Vamos encontrá-los e analisá-los ao vivo na plataforma do SIME.", h=0.95, qsize=13)
add_text(s, 0.7, 4.15, 8.0, 0.3, "Início da Demonstração Ao Vivo", size=11.5, bold=True, color=BRIGHT)
add_text(s, 0.7, 4.6, 9.6, 0.3, "➔ Abra o seu navegador e acesse a plataforma interna do SIME.", size=11, color="DCE4F2")
add_text(s, 0.7, 5.02, 9.6, 0.55, "➔ Tenha em mãos as dúvidas mais recorrentes de sua agência regional para investigarmos.", size=11, color="DCE4F2", spacing=1.1)
chip(s, 10.6, 4.6, 2.2, 0.55, "➔ Acesso ao Sistema", fill=BLUE, size=11)
add_text(s, 0.7, 7.1, 6.0, 0.22, "Capacitação SIME • Sebrae/PB", size=9, color="9FB0CC")
add_text(s, 12.1, 7.1, 0.83, 0.22, "Slide 12 / 17", size=8, color="9FB0CC", align=PP_ALIGN.RIGHT)

# ================================================================ SLIDE 13 — Interface
s = new_slide()
header(s, "Interface & Navegação", "Painel de Clientes Atendidos e Atendimentos")
card(s, 0.5, 1.62, 5.9, 5.05)
add_text(s, 0.8, 1.86, 5.0, 0.26, "ESTRUTURA DO PAINEL", size=9.5, bold=True, color=BLUE)
add_text(s, 0.8, 2.16, 5.0, 0.32, "Filtros e Visão Geral", size=13, bold=True, color=NAVY)
add_pic_cover(s, f"{MEDIA}/image9.png", 0.8, 2.6, 5.3, 1.72)
yy = 4.44
for t in ["Visão Geral: Total de clientes ➔ Total de atendimentos ➔ Perfil consolidado",
          "Filtros Principais: Período referencial ➔ Unidade/Agência ➔ Município",
          "Filtros Avançados: Porte empresarial ➔ Canal de atendimento ➔ Demais dimensões"]:
    add_text(s, 0.95, yy, 5.2, 0.4, [[("▪ ", {'bold': True, 'color': BLUE}), (t, {})]], size=9.2, color=SLATE, spacing=1.05)
    yy += 0.5
add_text(s, 0.8, 6.0, 5.3, 0.45, "* Utilize a combinação de múltiplos filtros para obter diagnósticos altamente específicos por território.",
         size=8.5, color=STEEL, spacing=1.05)
card(s, 6.93, 1.62, 5.9, 5.05)
add_text(s, 7.23, 1.86, 5.0, 0.26, "INTELIGÊNCIA DE GESTÃO", size=9.5, bold=True, color=BLUE)
add_text(s, 7.23, 2.16, 5.0, 0.32, "Perguntas que o painel responde", size=13, bold=True, color=NAVY)
add_pic_cover(s, f"{MEDIA}/image10.png", 7.23, 2.6, 5.3, 1.72)
yy = 4.44
for q in ["Quantos clientes individuais nós conseguimos atingir no período?",
          "Qual o volume absoluto de atendimentos prestados por nossa equipe?",
          "Qual o perfil de porte e localização desses clientes assistidos?",
          "Qual agência ou regional realizou a maior parcela desses esforços?"]:
    add_text(s, 7.4, yy, 5.1, 0.38, [[("? ", {'bold': True, 'color': BLUE, 'size': 11}), (q, {})]], size=9.2, color=SLATE, spacing=1.05)
    yy += 0.46
add_text(s, 7.23, 6.3, 5.3, 0.38, "Este painel serve para justificar o esforço operacional de atendimento da regional.",
         size=8.5, color=STEEL, spacing=1.05)
tip_bar(s, 0.5, 6.74, 12.33, "Dica:", "Certifique-se de limpar os filtros antes de iniciar uma nova pesquisa para evitar dados residuais!", h=0.4)
footer(s, num=13)

# ================================================================ SLIDE 14 — Acompanhamento
s = new_slide()
header(s, "Acompanhamento e Análise", "Painel de Metas Mobilizadoras")
card(s, 0.5, 1.62, 5.9, 5.05)
add_text(s, 0.8, 1.86, 5.0, 0.26, "VISUALIZAÇÃO ESTRATÉGICA", size=9.5, bold=True, color=BLUE)
add_text(s, 0.8, 2.16, 5.0, 0.32, "Componentes do Acompanhamento", size=13, bold=True, color=NAVY)
add_pic_cover(s, f"{MEDIA}/image11.png", 0.8, 2.6, 5.3, 1.72)
yy = 4.44
for c_ in ["Resultado Sebrae/PB: Desempenho geral consolidado do estado",
           "Meta Pactuada: O valor alvo formalmente acordado para o ano",
           "Percentual de Alcance: O quão próximo estamos de bater a meta de gestão",
           "Resultado por Unidade/Agência: Detalhamento tático para identificar gaps operacionais"]:
    add_text(s, 0.95, yy, 5.2, 0.4, [[("▪ ", {'bold': True, 'color': BLUE}), (c_, {})]], size=9.2, color=SLATE, spacing=1.05)
    yy += 0.48
add_text(s, 0.8, 6.34, 5.3, 0.32, "* Este painel é o principal instrumento de acompanhamento para as reuniões de resultados e tomadas de decisão estratégica das chefias.",
         size=8.5, color=STEEL, spacing=1.05)
card(s, 6.93, 1.62, 5.9, 5.05)
add_text(s, 7.23, 1.86, 5.0, 0.26, "PERGUNTA DE GESTÃO", size=9.5, bold=True, color=BLUE)
add_pic_cover(s, f"{MEDIA}/image12.png", 7.23, 2.24, 5.3, 2.0)
add_text(s, 7.23, 4.42, 5.3, 0.4, "“Estamos alcançando os resultados estratégicos que pactuamos?”", size=12.5, bold=True, color=NAVY, spacing=1.05)
add_text(s, 7.23, 4.95, 5.3, 0.9, [[("Foco em Solução: ", {'bold': True, 'color': BLUE}),
                                    ("Caso o percentual de alcance de uma regional esteja abaixo do planejado para o período, o painel auxilia a rastrear onde está o gargalo operacional.", {})]],
         size=9.2, color=SLATE, spacing=1.08)
tip_bar(s, 0.5, 6.74, 12.33, "Lembre-se:", "As Metas Mobilizadoras são calculadas com regras territoriais e não puramente operacionais!", h=0.4)
footer(s, num=14)

# ================================================================ SLIDE 15 — Desafio
s = new_slide()
header(s, "Desafio Prático Coletivo", "Exercício final — Da análise à ação")
card(s, 0.5, 1.62, 5.9, 5.0)
add_text(s, 0.8, 1.86, 5.0, 0.26, "SITUAÇÃO-PROBLEMA", size=9.5, bold=True, color=BLUE)
add_text(s, 0.8, 2.16, 5.0, 0.32, "Apresentação para a DIREX", size=13, bold=True, color=NAVY)
add_pic_cover(s, f"{MEDIA}/image13.png", 0.8, 2.6, 5.3, 1.6)
add_text(s, 0.8, 4.32, 5.3, 0.62,
         [[("Você foi convidado para apresentar, em uma reunião de Diretoria Executiva (", {}),
           ("DIREX", {'bold': True, 'color': NAVY}),
           ("), a situação atualizada de resultados operacionais e estratégicos da sua Agência/Unidade Regional.", {})]],
         size=9.2, color=SLATE, spacing=1.08)
add_text(s, 0.8, 4.96, 5.3, 0.7,
         [[("Utilizando o ", {}), ("SIME", {'bold': True, 'color': NAVY}),
           (", você deve analisar friamente os resultados atuais, mapear os pontos críticos e estruturar um plano de ação robusto e defensável para alcançar as metas estipuladas.", {})]],
         size=9.2, color=SLATE, spacing=1.08)
card(s, 0.8, 5.95, 5.3, 0.44, fill=TINT, line=None, radius=0.08)
add_text(s, 1.02, 6.06, 4.9, 0.26, "⏳ Tempo estimado para a atividade: 40 a 50 minutos.", size=9.5, bold=True, color=AMBER)
card(s, 6.93, 1.62, 5.9, 5.0)
add_text(s, 7.23, 1.86, 5.0, 0.26, "PLANO DE TRABALHO", size=9.5, bold=True, color=BLUE)
add_text(s, 7.23, 2.16, 5.0, 0.32, "O que cada grupo deve fazer", size=13, bold=True, color=NAVY)
add_pic_cover(s, f"{MEDIA}/image14.png", 7.23, 2.6, 5.3, 1.6)
yy = 4.32
steps15 = [
    [[("✔ ", {'bold': True, 'color': BLUE}), ("Acessar o ", {}), ("SIME", {'bold': True, 'color': NAVY}),
      (" e aplicar o filtro exclusivo da sua Agência/Unidade Regional.", {})]],
    [[("✔ ", {'bold': True, 'color': BLUE}), ("Identificar o percentual de alcance real de cada uma das Metas Mobilizadoras.", {})]],
    [[("✔ ", {'bold': True, 'color': BLUE}), ("Classificar as metas em: ", {}),
      ("dentro/acima do esperado", {'bold': True, 'color': NAVY}), (", em ", {}),
      ("atenção", {'bold': True, 'color': NAVY}), (", ou com ", {}),
      ("maior gap", {'bold': True, 'color': NAVY}), (".", {})]],
    [[("✔ ", {'bold': True, 'color': BLUE}), ("Cruzar com dados de clientes atendidos e atendimentos para buscar explicações técnicas.", {})]],
    [[("✔ ", {'bold': True, 'color': BLUE}), ("Escolher ", {}), ("até 3 prioridades estratégicas", {'bold': True, 'color': NAVY}),
      (" para evitar planos de ação genéricos.", {})]],
    [[("✔ ", {'bold': True, 'color': BLUE}), ("Construir ações táticas concretas direcionadas à resolução desses gargalos.", {})]],
]
for st_ in steps15:
    add_text(s, 7.4, yy, 5.15, 0.36, st_, size=8.8, color=SLATE, spacing=1.03)
    yy += 0.38
quote_bar(s, 0.5, 6.64, 12.33, "Não queremos uma apresentação de números. Queremos uma proposta de gestão.",
          tag="SIME ➔ Análise ➔ Decisão ➔ Ação ➔ Resultado", h=0.5, qsize=11.5)
footer(s, num=15)

# ================================================================ SLIDE 16 — Plano de Ação
s = new_slide()
header(s, "Instrumento de Gestão", "Plano de ação — do dado à entrega")
rows16 = [
    ["Meta / Indicador", "Result. Atual", "Gap", "O que os dados mostram?", "Ação proposta", "Expectativa de Resultado", "Prazo", "Responsável"],
    ["Cobertura MEI", "XX %", "XX pp", "Baixa cobertura observada em municípios específicos da regional.",
     "Realizar ação integrada de mobilização itinerante e mutirão de formalização de parceiros.",
     "Incorporar +150 novos MEIs à base da meta estadual.", "Set/26", "Agência"],
    ["Engajamento", "XX pts", "XX pts", "Alto volume de empresas com atendimento único, mas poucas recorrentes.",
     "Trabalhar ativação da carteira de empresas com maior potencial de trilhas de consultoria.",
     "Elevar índice de engajamento da carteira prioritária em 12%.", "Out/26", "Gestor"],
]
table_slide(s, 0.5, 1.66, 12.33, rows16, col_w=[1.7, 1.0, 0.85, 2.7, 2.75, 1.75, 0.78, 0.8],
            row_h=0.72, font_size=8.6, header_size=8.8)
tip_bar(s, 0.5, 4.05, 12.33, "📌 Importante:", "Cada ação proposta deve estar estritamente conectada ao gap identificado e estimar claramente quanto resultado pretende gerar.", h=0.46)
add_text(s, 0.5, 4.78, 12.33, 0.3, "Estrutura da Apresentação à DIREX (Apenas 5 minutos por grupo)", size=12, bold=True, color=NAVY)
steps16 = [
    ("1", "Onde estamos?", "Apresentação diagnóstica do estado atual."),
    ("2", "Onde chegar?", "Pactuações de metas planejadas."),
    ("3", "Principal problema?", "O maior gap ou ponto crítico mapeado."),
    ("4", "O que faremos?", "Ação estratégica e tática proposta."),
    ("5", "Quanto contribui?", "Expectativa numérica de entrega."),
]
sw, sgap = 2.33, 0.17
for i, (num, t, b) in enumerate(steps16):
    x = 0.5 + i * (sw + sgap)
    card(s, x, 5.16, sw, 1.62, fill=WHITE, line="E3E7EE", radius=0.07)
    add_rect(s, x, 5.16, sw, 0.07, fill=BLUE)
    add_text(s, x + 0.2, 5.34, 0.5, 0.4, num, size=20, bold=True, color=BLUE)
    add_text(s, x + 0.2, 5.88, sw - 0.4, 0.3, t, size=10.5, bold=True, color=NAVY, spacing=1.0)
    add_text(s, x + 0.2, 6.24, sw - 0.4, 0.45, b, size=8.5, color=SLATE, spacing=1.03)
footer(s, num=16)

# ================================================================ SLIDE 17 — Conclusão
s = new_slide()
header(s, "Conclusão e Próximos Passos", "O que quero que vocês levem daqui")
card(s, 0.5, 1.62, 5.9, 4.75)
add_text(s, 0.8, 1.86, 5.0, 0.26, "REGRAS E CONCEITOS", size=9.5, bold=True, color=BLUE)
add_pic_cover(s, f"{MEDIA}/image15.png", 0.8, 2.24, 5.3, 1.75)
yy = 4.18
for rg in ["Atendimento e cliente atendido são métricas com naturezas de contagem estruturalmente diferentes.",
           "Cada painel possui uma regra específica de cálculo e serve para responder a uma pergunta de negócio diferente.",
           "As Metas Mobilizadoras traduzem de forma prática a visão estratégica de longo prazo do Sebrae em resultados mensuráveis."]:
    add_text(s, 0.95, yy, 5.2, 0.5, [[("➔ ", {'bold': True, 'color': BLUE}), (rg, {})]], size=9.0, color=SLATE, spacing=1.05)
    yy += 0.64
card(s, 6.93, 1.62, 5.9, 4.75)
add_text(s, 7.23, 1.86, 5.0, 0.26, "GESTÃO ORIENTADA A DADOS", size=9.5, bold=True, color=BLUE)
add_pic_cover(s, f"{MEDIA}/image15.png", 7.23, 2.24, 5.3, 1.75)
yy = 4.18
for g in ["O SIME não serve apenas para acompanhar números ou preencher relatórios frios de metas.",
          "A plataforma serve fundamentalmente para identificar tendências e pontos onde precisamos intervir com plano de ação comercial.",
          "Tomadas de decisões táticas baseadas em relatórios do SIME mitigam desperdício de recursos regionais."]:
    add_text(s, 7.38, yy, 5.2, 0.5, [[("➔ ", {'bold': True, 'color': BLUE}), (g, {})]], size=9.0, color=SLATE, spacing=1.05)
    yy += 0.64
quote_bar(s, 0.5, 6.55, 12.33, "O dado mostra onde estamos. A gestão decide para onde vamos.", h=0.56, qsize=13)
footer(s, num=17)

prs.save(OUT)
print(f"OK: {OUT} com {len(prs.slides._sldIdLst)} slides")
