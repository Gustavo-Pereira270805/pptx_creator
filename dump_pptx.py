from pptx import Presentation
from pptx.util import Emu

prs = Presentation('original.pptx')

def in_(v):
    return round(Emu(v).inches, 3)

out = []
for idx, slide in enumerate(prs.slides, 1):
    out.append(f"\n{'='*90}\nSLIDE {idx}\n{'='*90}")
    for sh in slide.shapes:
        t = sh.shape_type
        name = sh.name
        pos = f"L={in_(sh.left)} T={in_(sh.top)} W={in_(sh.width)} H={in_(sh.height)}"
        kind = str(t)
        if sh.has_text_frame:
            txts = []
            for p in sh.text_frame.paragraphs:
                line = ''.join(r.text for r in p.runs)
                if line.strip():
                    txts.append(line)
            if txts:
                out.append(f"\n[{kind}] '{name}' {pos}")
                for ln in txts:
                    out.append(f"    | {ln}")
        elif sh.has_table:
            out.append(f"\n[{kind}] TABLE '{name}' {pos}")
            for r in sh.table.rows:
                out.append(f"    | {' | '.join(c.text for c in r.cells)}")
        else:
            out.append(f"\n[{kind}] '{name}' {pos}  (no text)")

with open('/tmp/dump.txt', 'w') as f:
    f.write('\n'.join(out))
print("written", len(out), "lines")
